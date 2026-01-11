from fastapi import FastAPI, File, UploadFile, BackgroundTasks
from fastapi.responses import FileResponse
from pdf2docx import Converter
import os
import io                                     
from pypdf import PdfReader, PdfWriter
from fastapi.responses import StreamingResponse
import pandas as pd
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt

app = FastAPI()

# دالة كتمسح الملف من بعد ما يوصل للمستخدم
def remove_file(path: str):
    try:
        os.remove(path)
    except Exception:
        pass

@app.post("/convert-to-word")
async def convert_pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    pdf_filename = file.filename
    docx_filename = pdf_filename.replace(".pdf", ".docx")

    # حفظ PDF
    with open(pdf_filename, "wb") as buffer:
        buffer.write(await file.read())

    # التحويل
    try:
        cv = Converter(pdf_filename)
        cv.convert(docx_filename)
        cv.close()
    except Exception as e:
        return {"error": str(e)}

    # جدولة مسح الملفات بعد الإرسال
    background_tasks.add_task(remove_file, pdf_filename)
    background_tasks.add_task(remove_file, docx_filename)

    return FileResponse(docx_filename, filename=docx_filename)


@app.post("/compress-pdf")
async def compress_pdf(file: UploadFile = File(...)):
    try:
        # 1. قراءة الملف
        reader = PdfReader(file.file)
        writer = PdfWriter()

        # 2. نسخ الصفحات وضغطها
        for page in reader.pages:
            writer.add_page(page)
        
        # تفعيل خوارزميات الضغط
        writer.compress_identical_objects = True 

        # 3. الحفظ في الذاكرة
        output_stream = io.BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)

        # 4. إرسال النتيجة
        return StreamingResponse(
            output_stream, 
            media_type="application/pdf", 
            headers={"Content-Disposition": "attachment; filename=compressed.pdf"}
        )
    except Exception as e:
        return {"error": str(e)}
    
@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    try:
        # Save uploaded file temporarily
        temp_pdf = f"temp_{file.filename}"
        with open(temp_pdf, "wb") as f:
            f.write(await file.read())

        excel_filename = "converted_tables.xlsx"
        
        # Open PDF and extract tables
        with pdfplumber.open(temp_pdf) as pdf:
            with pd.ExcelWriter(excel_filename, engine='openpyxl') as writer:
                has_tables = False
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        # Convert to DataFrame
                        df = pd.DataFrame(table[1:], columns=table[0])
                        # Save to Excel Sheet
                        sheet_name = f"Page{i+1}_Table{j+1}"
                        df.to_excel(writer, sheet_name=sheet_name, index=False)
                        has_tables = True
                
                # If no tables found, create an empty sheet
                if not has_tables:
                    pd.DataFrame(["No tables found"]).to_excel(writer, sheet_name="Info")

        # Clean up PDF
        os.remove(temp_pdf)

        # Return the Excel file
        return FileResponse(
            excel_filename, 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
            filename="tables.xlsx"
        )
    except Exception as e:
        return {"error": str(e)}    
    
@app.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    try:
        temp_pdf = f"temp_{file.filename}"
        with open(temp_pdf, "wb") as f:
            f.write(await file.read())

        ppt_filename = "converted_slides.pptx"
        prs = Presentation()
        
        # Simple Text Extraction Strategy
        with pdfplumber.open(temp_pdf) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                # Create blank slide
                slide_layout = prs.slide_layouts[6] 
                slide = prs.slides.add_slide(slide_layout)
                
                # Add text box
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                tf = txBox.text_frame
                tf.text = text if text else "Image Slide (Content not extracting)"

        prs.save(ppt_filename)
        os.remove(temp_pdf)

        return FileResponse(
            ppt_filename, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename="presentation.pptx"
        )
    except Exception as e:
        return {"error": str(e)}    